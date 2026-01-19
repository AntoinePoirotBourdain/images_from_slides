#!/usr/bin/env python3
"""
Script pour extraire les slides d'un PPTX en PNG avec fond transparent
Le titre est extrait puis supprimé de l'image
Les images sont croppées au contenu minimal et compressées
Les images existantes avec le même nom sont écrasées
"""

import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw, ImageChops
import subprocess
import tempfile
import shutil

def sanitize_filename(text):
    """Nettoie le texte pour en faire un nom de fichier valide"""
    invalid_chars = '<>:"/\\|?*'
    for char in invalid_chars:
        text = text.replace(char, '')
    
    text = ' '.join(text.split())
    
    if len(text) > 100:
        text = text[:100]
    
    return text.strip()

def extract_slide_title_and_bbox(slide):
    """
    Extrait le titre d'une slide et sa bounding box (position)
    Retourne (titre, (x, y, width, height)) en coordonnées relatives
    """
    # Récupérer les dimensions de la présentation
    prs = slide.part.package.presentation_part.presentation
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title = title_shape.text
        
        # Position et dimensions du titre (en EMU)
        left = title_shape.left
        top = title_shape.top
        width = title_shape.width
        height = title_shape.height
        
        # Convertir en coordonnées relatives (0-1)
        bbox = (
            left / slide_width,
            top / slide_height,
            width / slide_width,
            height / slide_height
        )
        
        return title, bbox
    
    # Si pas de titre défini, chercher le premier textbox
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            title = shape.text.strip()
            
            bbox = (
                shape.left / slide_width,
                shape.top / slide_height,
                shape.width / slide_width,
                shape.height / slide_height
            )
            
            return title, bbox
    
    return None, None

def remove_title_from_image(img, bbox, margin=10):
    """
    Supprime le titre de l'image en rendant cette zone transparente
    
    Args:
        img: Image PIL
        bbox: (x_rel, y_rel, width_rel, height_rel) coordonnées relatives
        margin: Marge supplémentaire en pixels autour du titre
    """
    if bbox is None:
        return img
    
    img_width, img_height = img.size
    x_rel, y_rel, width_rel, height_rel = bbox
    
    # Convertir en coordonnées pixels
    x = int(x_rel * img_width) - margin
    y = int(y_rel * img_height) - margin
    width = int(width_rel * img_width) + 2 * margin
    height = int(height_rel * img_height) + 2 * margin
    
    # S'assurer que les coordonnées sont dans les limites
    x = max(0, x)
    y = max(0, y)
    width = min(width, img_width - x)
    height = min(height, img_height - y)
    
    # Créer un masque pour rendre transparent
    if img.mode != 'RGBA':
        img = img.convert('RGBA')
    
    # Dessiner un rectangle transparent
    draw = ImageDraw.Draw(img)
    draw.rectangle([x, y, x + width, y + height], fill=(255, 255, 255, 0))
    
    return img

def autocrop_image(img, margin=0):
    """
    Crop automatiquement l'image à son contenu non-transparent
    
    Args:
        img: Image PIL en mode RGBA
        margin: Marge en pixels à garder autour du contenu
    
    Returns:
        Image croppée
    """
    if img.mode != 'RGBA':
        img = img.convert('RGBA')
    
    # Obtenir le canal alpha
    alpha = img.split()[-1]
    
    # Trouver la bounding box du contenu non-transparent
    bbox = alpha.getbbox()
    
    if bbox is None:
        # Image entièrement transparente, retourner une petite image
        return img.crop((0, 0, 100, 100))
    
    # Ajouter une marge
    left, top, right, bottom = bbox
    
    left = max(0, left - margin)
    top = max(0, top - margin)
    right = min(img.width, right + margin)
    bottom = min(img.height, bottom + margin)
    
    # Crop l'image
    return img.crop((left, top, right, bottom))

def resize_image(img, scale_percent):
    """
    Redimensionne l'image selon un pourcentage
    
    Args:
        img: Image PIL
        scale_percent: Pourcentage de la taille originale (ex: 20 pour 20%)
    
    Returns:
        Image redimensionnée
    """
    if scale_percent >= 100:
        return img
    
    width = int(img.width * scale_percent / 100)
    height = int(img.height * scale_percent / 100)
    
    # Utiliser LANCZOS pour une meilleure qualité lors de la réduction
    return img.resize((width, height), Image.Resampling.LANCZOS)

def get_file_size_mb(filepath):
    """Retourne la taille du fichier en MB"""
    return os.path.getsize(filepath) / (1024 * 1024)

def pptx_to_png_transparent(pptx_path, output_dir='images', dpi=300, remove_title=True, 
                           autocrop=True, crop_margin=20, scale_percent=100):
    """
    Convertit un PPTX en PNG avec fond transparent et titre supprimé
    
    Args:
        pptx_path: Chemin vers le fichier PPTX
        output_dir: Dossier de sortie pour les images
        dpi: Résolution des images (300 recommandé)
        remove_title: Si True, supprime le titre de l'image
        autocrop: Si True, crop l'image au contenu minimal
        crop_margin: Marge en pixels à garder autour du contenu lors du crop
        scale_percent: Pourcentage de réduction de la taille (ex: 20 pour 20% de la taille originale)
    """
    
    if not os.path.exists(pptx_path):
        print(f"❌ Fichier non trouvé : {pptx_path}")
        return
    
    # Créer le dossier de sortie s'il n'existe pas
    os.makedirs(output_dir, exist_ok=True)
    
    # Compter les fichiers existants
    existing_files = [f for f in os.listdir(output_dir) if f.lower().endswith('.png')]
    if existing_files:
        print(f"📁 Dossier de sortie : {output_dir}/ ({len(existing_files)} images existantes)")
        print(f"⚠️  Les images avec le même nom seront écrasées")
    else:
        print(f"📁 Dossier de sortie : {output_dir}/")
    
    print(f"📂 Chargement de {pptx_path}...")
    
    try:
        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)
        
        print(f"📊 {total_slides} slides trouvées")
        print(f"🎨 Export en PNG avec fond transparent (DPI: {dpi})")
        if remove_title:
            print(f"✂️  Les titres seront supprimés des images")
        if autocrop:
            print(f"✂️  Les images seront croppées au contenu minimal (marge: {crop_margin}px)")
        if scale_percent < 100:
            print(f"📉 Les images seront réduites à {scale_percent}% de leur taille")
        print()
        
        # Extraire les titres et bboxes AVANT la conversion
        slides_info = []
        for idx, slide in enumerate(prs.slides, start=1):
            title, bbox = extract_slide_title_and_bbox(slide)
            slides_info.append((title, bbox))
            
            if title:
                print(f"  📄 Slide {idx}: {title[:50]}{'...' if len(title) > 50 else ''}")
            else:
                print(f"  📄 Slide {idx}: (sans titre)")
        
        print()
        
        # Vérifier si LibreOffice est installé
        has_libreoffice = shutil.which('soffice') or shutil.which('libreoffice')
        
        if has_libreoffice:
            print("✅ LibreOffice détecté - Utilisation pour meilleure qualité")
            use_libreoffice_method(pptx_path, output_dir, slides_info, dpi, remove_title, 
                                 autocrop, crop_margin, scale_percent)
        else:
            print("❌ LibreOffice NON détecté - OBLIGATOIRE pour cette fonctionnalité")
            print("   Installez LibreOffice:")
            print("   - Ubuntu/Debian: sudo apt-get install libreoffice")
            print("   - macOS: brew install --cask libreoffice")
            print("   - Windows: https://www.libreoffice.org/download/")
            return
        
        print()
        print(f"✅ Export terminé ! Images dans {output_dir}/")
        
    except Exception as e:
        print(f"❌ Erreur : {e}")
        import traceback
        traceback.print_exc()

def use_libreoffice_method(pptx_path, output_dir, slides_info, dpi, remove_title, 
                          autocrop, crop_margin, scale_percent):
    """Méthode avec LibreOffice"""
    
    with tempfile.TemporaryDirectory() as tmpdir:
        # Convertir PPTX en PDF
        print("🔄 Conversion PPTX → PDF...")
        cmd = [
            'soffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', tmpdir,
            pptx_path
        ]
        
        try:
            result = subprocess.run(cmd, check=True, capture_output=True, text=True)
        except FileNotFoundError:
            cmd[0] = 'libreoffice'
            result = subprocess.run(cmd, check=True, capture_output=True, text=True)
        except subprocess.CalledProcessError as e:
            print(f"❌ Erreur lors de la conversion PDF : {e}")
            print(f"Sortie : {e.output}")
            return
        
        pdf_file = os.path.join(tmpdir, os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf')
        
        if not os.path.exists(pdf_file):
            print(f"❌ Fichier PDF non créé : {pdf_file}")
            return
        
        # Convertir PDF en PNG
        print("🔄 Conversion PDF → PNG...")
        try:
            subprocess.run([
                'pdftoppm',
                '-png',
                '-r', str(dpi),
                pdf_file,
                os.path.join(tmpdir, 'slide')
            ], check=True)
        except FileNotFoundError:
            print("❌ pdftoppm non trouvé. Installez poppler-utils:")
            print("   - Ubuntu/Debian: sudo apt-get install poppler-utils")
            print("   - macOS: brew install poppler")
            return
        
        png_files = sorted([f for f in os.listdir(tmpdir) if f.endswith('.png')])
        
        if not png_files:
            print(f"❌ Aucune image PNG générée dans {tmpdir}")
            return
        
        print("🔄 Traitement des images...")
        
        # Dictionnaire pour gérer les noms en double dans la présentation
        name_counts = {}
        total_size = 0
        files_overwritten = 0
        files_created = 0
        
        for idx, (png_file, (title, bbox)) in enumerate(zip(png_files, slides_info), start=1):
            # Définir le nom de sortie (SANS numéro)
            if title:
                safe_title = sanitize_filename(title)
                
                # Gérer les doublons dans la présentation
                if safe_title in name_counts:
                    name_counts[safe_title] += 1
                    output_name = f"{safe_title}_{name_counts[safe_title]}.png"
                else:
                    name_counts[safe_title] = 1
                    output_name = f"{safe_title}.png"
            else:
                output_name = f"Slide_{idx}.png"
            
            output_path = os.path.join(output_dir, output_name)
            
            # Vérifier si on écrase un fichier existant
            file_existed = os.path.exists(output_path)
            
            # Charger l'image
            img_path = os.path.join(tmpdir, png_file)
            img = Image.open(img_path)
            
            # Récupérer la taille originale
            original_size = img.size
            
            # Convertir en RGBA
            if img.mode != 'RGBA':
                img = img.convert('RGBA')
            
            # Rendre le fond blanc transparent
            data = img.getdata()
            new_data = []
            
            for item in data:
                # Si le pixel est blanc (ou presque), le rendre transparent
                if item[0] > 250 and item[1] > 250 and item[2] > 250:
                    new_data.append((255, 255, 255, 0))
                else:
                    new_data.append(item)
            
            img.putdata(new_data)
            
            # Supprimer le titre si demandé
            if remove_title and bbox:
                img = remove_title_from_image(img, bbox, margin=15)
            
            # Crop automatique si demandé
            if autocrop:
                img = autocrop_image(img, margin=crop_margin)
            
            # Réduire la taille si demandé
            if scale_percent < 100:
                img = resize_image(img, scale_percent)
            
            # Sauvegarder avec optimisation PNG (écrase si existe)
            img.save(output_path, 'PNG', optimize=True)
            
            # Calculer la taille du fichier
            file_size = get_file_size_mb(output_path)
            total_size += file_size
            
            # Compter les fichiers créés vs écrasés
            if file_existed:
                files_overwritten += 1
                status_icon = "♻️"
            else:
                files_created += 1
                status_icon = "✓"
            
            # Afficher les infos
            status_parts = []
            if remove_title and bbox:
                status_parts.append("✂️ titre")
            if autocrop:
                status_parts.append(f"📐 {img.size[0]}x{img.size[1]}")
            if scale_percent < 100:
                status_parts.append(f"📉 {scale_percent}%")
            status_parts.append(f"💾 {file_size:.2f}MB")
            
            status = " | ".join(status_parts)
            print(f"  {status_icon} {status} → {output_name}")
        
        print(f"\n📊 Statistiques:")
        if files_created > 0:
            print(f"   ✅ {files_created} nouvelles images créées")
        if files_overwritten > 0:
            print(f"   ♻️  {files_overwritten} images écrasées")
        print(f"   💾 Taille totale : {total_size:.2f} MB")
        print(f"   💾 Taille moyenne : {total_size/len(png_files):.2f} MB par image")

# Utilisation principale
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python extract_slides.py <fichier.pptx> [dossier_sortie] [dpi] [options]")
        print()
        print("Options:")
        print("  --keep-title           Garder le titre dans l'image")
        print("  --no-crop              Ne pas cropper l'image")
        print("  --crop-margin=N        Marge de crop en pixels (défaut: 20)")
        print("  --scale=N              Échelle en pourcentage (défaut: 100)")
        print()
        print("Exemples:")
        print("  python extract_slides.py presentation.pptx")
        print("  python extract_slides.py presentation.pptx images")
        print("  python extract_slides.py presentation.pptx images 300")
        print("  python extract_slides.py presentation.pptx images 300 --scale=50")
        print("  python extract_slides.py presentation.pptx images 300 --scale=20")
        print("  python extract_slides.py presentation.pptx images 300 --keep-title --no-crop")
        print()
        print("Par défaut:")
        print("  - Dossier de sortie: images/")
        print("  - DPI: 300")
        print("  - Échelle: 100% (taille originale après crop)")
        print("  - Les titres sont supprimés des images")
        print("  - Les images sont croppées avec marge de 20px")
        print("  - Les images existantes avec le même nom sont écrasées")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    output_folder = sys.argv[2] if len(sys.argv) > 2 else 'images'
    
    # Vérifier les options
    resolution = 300
    keep_title = False
    do_autocrop = True
    margin = 20
    scale = 100  # Par défaut 100% (taille originale)
    
    for arg in sys.argv[3:]:
        if arg == '--keep-title':
            keep_title = True
        elif arg == '--no-crop':
            do_autocrop = False
        elif arg.startswith('--crop-margin='):
            margin = int(arg.split('=')[1])
        elif arg.startswith('--scale='):
            scale = int(arg.split('=')[1])
        elif arg.isdigit():
            resolution = int(arg)
    
    pptx_to_png_transparent(
        pptx_file, 
        output_folder, 
        resolution,
        remove_title=not keep_title,
        autocrop=do_autocrop,
        crop_margin=margin,
        scale_percent=scale
    )